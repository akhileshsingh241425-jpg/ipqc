"""
═══════════════════════════════════════════════════════════════
QMS Document Search Engine
Advanced text extraction, indexing & intelligent search
Solar Panel Manufacturing QMS | ISO 9001:2015
═══════════════════════════════════════════════════════════════
"""

import os
import re
import math
import json
import logging
from collections import Counter, defaultdict

logger = logging.getLogger(__name__)

# ═══════════════════════════════════════════════
# TEXT EXTRACTION
# ═══════════════════════════════════════════════

def extract_text_from_file(file_path):
    """Extract text from any supported file type"""
    if not file_path or not os.path.exists(file_path):
        return ''
    
    ext = file_path.rsplit('.', 1)[-1].lower() if '.' in file_path else ''
    
    extractors = {
        'pdf': _extract_pdf,
        'docx': _extract_docx,
        'doc': _extract_docx,
        'xlsx': _extract_xlsx,
        'xls': _extract_xlsx,
        'txt': _extract_text,
        'csv': _extract_text,
        'pptx': _extract_pptx,
    }
    
    extractor = extractors.get(ext)
    if not extractor:
        return ''
    
    try:
        text = extractor(file_path)
        # Clean up extracted text
        text = _clean_text(text)
        return text
    except Exception as e:
        logger.error(f"Text extraction failed for {file_path}: {e}")
        return ''


def _extract_pdf(file_path):
    """Extract text from PDF using PyPDF2"""
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(file_path)
        text_parts = []
        for page_num, page in enumerate(reader.pages, 1):
            page_text = page.extract_text()
            if page_text:
                text_parts.append(f"[Page {page_num}]\n{page_text}")
        return '\n\n'.join(text_parts)
    except Exception as e:
        logger.error(f"PDF extraction error: {e}")
        return ''


def _extract_docx(file_path):
    """Extract text from DOCX"""
    try:
        from docx import Document as DocxDocument
        doc = DocxDocument(file_path)
        text_parts = []
        
        # Extract paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                # Include heading level info
                if para.style and para.style.name and 'Heading' in para.style.name:
                    text_parts.append(f"\n## {para.text}\n")
                else:
                    text_parts.append(para.text)
        
        # Extract tables
        for table_idx, table in enumerate(doc.tables, 1):
            text_parts.append(f"\n[Table {table_idx}]")
            for row in table.rows:
                row_text = ' | '.join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    text_parts.append(row_text)
        
        return '\n'.join(text_parts)
    except ImportError:
        logger.warning("python-docx not installed, DOCX extraction unavailable")
        return ''
    except Exception as e:
        logger.error(f"DOCX extraction error: {e}")
        return ''


def _extract_xlsx(file_path):
    """Extract text from Excel files"""
    try:
        from openpyxl import load_workbook
        wb = load_workbook(file_path, read_only=True, data_only=True)
        text_parts = []
        
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            text_parts.append(f"\n[Sheet: {sheet_name}]")
            for row in ws.iter_rows(values_only=True):
                row_text = ' | '.join(str(cell) for cell in row if cell is not None)
                if row_text.strip():
                    text_parts.append(row_text)
        
        wb.close()
        return '\n'.join(text_parts)
    except Exception as e:
        logger.error(f"Excel extraction error: {e}")
        return ''


def _extract_pptx(file_path):
    """Extract text from PowerPoint files"""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            text_parts.append(f"\n[Slide {slide_num}]")
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    text_parts.append(shape.text)
        
        return '\n'.join(text_parts)
    except ImportError:
        logger.warning("python-pptx not installed")
        return ''
    except Exception as e:
        logger.error(f"PPTX extraction error: {e}")
        return ''


def _extract_text(file_path):
    """Extract text from plain text/CSV files"""
    try:
        encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
        for enc in encodings:
            try:
                with open(file_path, 'r', encoding=enc) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue
        return ''
    except Exception as e:
        logger.error(f"Text extraction error: {e}")
        return ''


def _clean_text(text):
    """Clean and normalize extracted text"""
    if not text:
        return ''
    # Remove excessive whitespace but keep paragraph structure
    text = re.sub(r'[ \t]+', ' ', text)
    text = re.sub(r'\n{4,}', '\n\n\n', text)
    # Remove null bytes
    text = text.replace('\x00', '')
    return text.strip()


# ═══════════════════════════════════════════════
# TEXT CHUNKING
# ═══════════════════════════════════════════════

def chunk_text(text, chunk_size=500, overlap=100):
    """Break text into overlapping chunks for better search"""
    if not text:
        return []
    
    # Split into paragraphs first
    paragraphs = re.split(r'\n\s*\n', text)
    chunks = []
    current_chunk = ''
    
    for para in paragraphs:
        para = para.strip()
        if not para:
            continue
        
        if len(current_chunk) + len(para) < chunk_size:
            current_chunk += ('\n\n' if current_chunk else '') + para
        else:
            if current_chunk:
                chunks.append(current_chunk.strip())
            
            # If single paragraph is larger than chunk_size, split by sentences
            if len(para) > chunk_size:
                sentences = re.split(r'(?<=[.!?])\s+', para)
                current_chunk = ''
                for sent in sentences:
                    if len(current_chunk) + len(sent) < chunk_size:
                        current_chunk += (' ' if current_chunk else '') + sent
                    else:
                        if current_chunk:
                            chunks.append(current_chunk.strip())
                        current_chunk = sent
            else:
                current_chunk = para
    
    if current_chunk:
        chunks.append(current_chunk.strip())
    
    return chunks


# ═══════════════════════════════════════════════
# SEARCH ENGINE
# ═══════════════════════════════════════════════

# Stopwords (English + Hindi common words)
STOPWORDS = {
    'the', 'a', 'an', 'is', 'are', 'was', 'were', 'be', 'been', 'being',
    'have', 'has', 'had', 'do', 'does', 'did', 'will', 'would', 'could',
    'should', 'may', 'might', 'shall', 'can', 'of', 'in', 'to', 'for',
    'with', 'on', 'at', 'by', 'from', 'as', 'into', 'through', 'during',
    'before', 'after', 'above', 'below', 'between', 'and', 'but', 'or',
    'not', 'no', 'so', 'if', 'this', 'that', 'these', 'those', 'it',
    'its', 'he', 'she', 'they', 'we', 'i', 'you', 'me', 'him', 'her',
    'us', 'them', 'my', 'your', 'his', 'our', 'their', 'what', 'which',
    'who', 'when', 'where', 'how', 'why', 'all', 'each', 'every', 'both',
    'few', 'more', 'most', 'than', 'very', 'just', 'also', 'only',
    # Hindi common
    'ka', 'ke', 'ki', 'ko', 'me', 'hai', 'hain', 'ye', 'wo', 'kya',
    'kaise', 'kab', 'kaha', 'se', 'par', 'aur', 'ya', 'tho', 'to',
    'mein', 'mujhe', 'koi', 'bhi', 'nahi', 'wala', 'wali', 'wale',
    'chahiye', 'karo', 'kare', 'batao', 'bata', 'de', 'do',
}

# QMS Synonym expansion for better matching
QMS_SYNONYMS = {
    'sop': ['procedure', 'standard operating procedure', 'work instruction'],
    'wi': ['work instruction', 'procedure', 'method'],
    'ncr': ['non conformance', 'non-conformance', 'nonconformity', 'defect', 'rejection'],
    'capa': ['corrective action', 'preventive action', 'correction', 'improvement'],
    'calibration': ['calibrate', 'calibrated', 'cal', 'measurement', 'instrument'],
    'inspection': ['inspect', 'check', 'verification', 'testing', 'examination'],
    'lamination': ['laminate', 'laminator', 'laminated'],
    'soldering': ['solder', 'stringer', 'stringing', 'ribbon'],
    'packing': ['packaging', 'pack', 'dispatch', 'shipping'],
    'testing': ['test', 'flash', 'el', 'hipot', 'hi-pot', 'iv curve'],
    'framing': ['frame', 'framing machine', 'corner key', 'sealant'],
    'layup': ['lay-up', 'glass', 'eva', 'epe', 'backsheet'],
    'junction box': ['jbox', 'j-box', 'junction', 'potting', 'diode'],
    'pdi': ['pre delivery', 'pre-delivery', 'final inspection'],
    'bis': ['bureau of indian standards', 'isi', 'certification'],
    'almm': ['approved list', 'module manufacturer'],
    'iec': ['international electrotechnical', 'standard'],
    'temperature': ['temp', 'temperature profile', 'thermocouple'],
    'iso': ['iso 9001', 'quality management system', 'qms'],
    'audit': ['auditing', 'assessment', 'review', 'surveillance'],
    'training': ['train', 'competency', 'skill', 'qualification'],
    'supplier': ['vendor', 'supply chain', 'procurement'],
    'specification': ['spec', 'specs', 'requirement', 'standard'],
    'document': ['doc', 'documentation', 'record', 'file'],
    'quality': ['quality control', 'qc', 'qa', 'quality assurance'],
    'production': ['manufacturing', 'production line', 'assembly'],
    'material': ['raw material', 'bom', 'bill of material', 'component'],
    'cell': ['solar cell', 'photovoltaic', 'pv cell'],
    'module': ['solar module', 'panel', 'solar panel', 'pv module'],
    'rejection': ['reject', 'rejected', 'defective', 'rework'],
    'approval': ['approve', 'authorized', 'sanction', 'signed'],
    'revision': ['revise', 'revised', 'version', 'amendment'],
    'controlled': ['controlled copy', 'master copy', 'original'],
}


def tokenize(text):
    """Tokenize text into meaningful words"""
    if not text:
        return []
    text = text.lower()
    # Keep alphanumeric + hyphen
    tokens = re.findall(r'[a-z0-9][\w\-]*[a-z0-9]|[a-z0-9]', text)
    return [t for t in tokens if t not in STOPWORDS and len(t) > 1]


def expand_query(query):
    """Expand query with synonyms for better recall"""
    tokens = tokenize(query)
    expanded = set(tokens)
    
    query_lower = query.lower()
    for key, syns in QMS_SYNONYMS.items():
        if key in query_lower or any(t in key.split() for t in tokens):
            for syn in syns:
                expanded.update(tokenize(syn))
        for syn in syns:
            if syn in query_lower:
                expanded.update(tokenize(key))
                break
    
    return list(expanded)


def compute_tf(tokens):
    """Compute term frequency"""
    counter = Counter(tokens)
    total = len(tokens) if tokens else 1
    return {t: c / total for t, c in counter.items()}


def search_documents(query, documents, top_k=10):
    """
    Advanced document search with TF-IDF scoring.
    
    Args:
        query: Search query string
        documents: List of dicts with 'id', 'title', 'extracted_text', 'description', etc.
        top_k: Number of top results to return
    
    Returns:
        List of result dicts with score, passages, document info
    """
    if not query or not documents:
        return []
    
    query_tokens = expand_query(query)
    if not query_tokens:
        return []
    
    # Build document corpus
    doc_chunks = []  # (doc_info, chunk_text, chunk_idx)
    for doc in documents:
        text = doc.get('extracted_text', '') or ''
        if not text:
            # Fall back to description + title
            text = f"{doc.get('title', '')} {doc.get('description', '')}"
        
        chunks = chunk_text(text) if len(text) > 500 else [text]
        for idx, chunk in enumerate(chunks):
            doc_chunks.append((doc, chunk, idx))
    
    if not doc_chunks:
        return []
    
    # Compute IDF across all chunks
    N = len(doc_chunks)
    df = Counter()
    chunk_tokens_list = []
    for _, chunk, _ in doc_chunks:
        tokens = set(tokenize(chunk))
        chunk_tokens_list.append(tokens)
        for t in tokens:
            df[t] += 1
    
    idf = {}
    for t in df:
        idf[t] = math.log(1 + N / (1 + df[t]))
    
    # Score each chunk
    results = []
    for i, (doc, chunk, chunk_idx) in enumerate(doc_chunks):
        chunk_token_list = tokenize(chunk)
        tf = compute_tf(chunk_token_list)
        chunk_token_set = chunk_tokens_list[i]
        
        # TF-IDF score
        score = 0
        matched_terms = []
        for qt in query_tokens:
            if qt in chunk_token_set:
                tfidf = tf.get(qt, 0) * idf.get(qt, 1)
                score += tfidf
                matched_terms.append(qt)
        
        if score <= 0:
            continue
        
        # Boost for exact phrase match
        query_lower = query.lower()
        chunk_lower = chunk.lower()
        if query_lower in chunk_lower:
            score *= 2.5
        
        # Boost for title match
        title_lower = (doc.get('title', '') or '').lower()
        if any(qt in tokenize(title_lower) for qt in query_tokens):
            score *= 1.5
        
        # Boost for more matched terms (coverage)
        coverage = len(matched_terms) / max(len(query_tokens), 1)
        score *= (1 + coverage)
        
        # Extract the most relevant passage (highlighted)
        passage = _extract_passage(chunk, query_tokens, max_length=350)
        
        results.append({
            'doc_id': doc.get('id'),
            'doc_number': doc.get('doc_number', ''),
            'title': doc.get('title', ''),
            'category': doc.get('category', ''),
            'department': doc.get('department', ''),
            'status': doc.get('status', ''),
            'file_name': doc.get('file_name', ''),
            'score': round(score, 4),
            'passage': passage,
            'chunk_index': chunk_idx,
            'matched_terms': list(set(matched_terms)),
            'coverage': round(coverage * 100, 1),
        })
    
    # Sort by score descending, deduplicate by doc_id (keep top passage per doc)
    results.sort(key=lambda x: x['score'], reverse=True)
    
    # Keep top passages but also group by document
    seen_docs = {}
    final_results = []
    for r in results:
        doc_id = r['doc_id']
        if doc_id not in seen_docs:
            seen_docs[doc_id] = r
            r['additional_passages'] = []
            final_results.append(r)
        else:
            # Add as additional passage
            if len(seen_docs[doc_id]['additional_passages']) < 3:
                seen_docs[doc_id]['additional_passages'].append({
                    'passage': r['passage'],
                    'score': r['score'],
                    'chunk_index': r['chunk_index']
                })
    
    return final_results[:top_k]


def _extract_passage(text, query_tokens, max_length=350):
    """Extract the most relevant passage around matched terms"""
    if not text:
        return ''
    
    text_lower = text.lower()
    
    # Find the best position (most query term matches nearby)
    sentences = re.split(r'(?<=[.!?\n])\s+', text)
    if not sentences:
        return text[:max_length]
    
    best_score = -1
    best_idx = 0
    
    for i, sent in enumerate(sentences):
        sent_lower = sent.lower()
        score = sum(1 for qt in query_tokens if qt in sent_lower)
        if score > best_score:
            best_score = score
            best_idx = i
    
    # Build passage from best sentence + context
    start = max(0, best_idx - 1)
    end = min(len(sentences), best_idx + 3)
    passage = ' '.join(sentences[start:end])
    
    if len(passage) > max_length:
        passage = passage[:max_length] + '...'
    
    # Highlight matched terms
    for qt in query_tokens:
        pattern = re.compile(re.escape(qt), re.IGNORECASE)
        passage = pattern.sub(f'**{qt}**', passage)
    
    return passage


# ═══════════════════════════════════════════════
# QUESTION ANSWERING
# ═══════════════════════════════════════════════

def answer_question(query, search_results, all_docs_count=0, indexed_count=0):
    """
    Generate a comprehensive answer combining search results with context.
    Returns a structured response dict.
    """
    query_lower = query.lower()
    
    response = {
        'answer': '',
        'sources': [],
        'total_docs': all_docs_count,
        'indexed_docs': indexed_count,
        'results_count': len(search_results),
        'confidence': 'low',
        'suggestions': []
    }
    
    if not search_results:
        # No results found
        response['answer'] = _generate_no_result_answer(query_lower)
        response['suggestions'] = _generate_suggestions(query_lower)
        return response
    
    # Build answer from top results
    top_result = search_results[0]
    
    # Determine confidence
    if top_result['score'] > 1.0 and top_result['coverage'] > 60:
        response['confidence'] = 'high'
    elif top_result['score'] > 0.3:
        response['confidence'] = 'medium'
    else:
        response['confidence'] = 'low'
    
    # Build answer text
    answer_parts = []
    
    if response['confidence'] == 'high':
        answer_parts.append(f"📄 **\"{top_result['title']}\"** ({top_result['doc_number']}) mein ye information mili:\n")
    else:
        answer_parts.append(f"🔍 Related information {len(search_results)} document(s) mein mili:\n")
    
    # Add top passages
    for i, result in enumerate(search_results[:5]):
        passage = result['passage']
        if passage:
            source_label = f"**{result['title']}**"
            if result.get('doc_number'):
                source_label += f" ({result['doc_number']})"
            
            answer_parts.append(f"\n📌 {source_label}:")
            answer_parts.append(f"_{passage}_\n")
            
            # Add additional passages from same doc
            for ap in result.get('additional_passages', [])[:1]:
                if ap['passage']:
                    answer_parts.append(f"  _{ap['passage']}_\n")
        
        response['sources'].append({
            'id': result['doc_id'],
            'title': result['title'],
            'doc_number': result['doc_number'],
            'category': result['category'],
            'score': result['score'],
            'coverage': result['coverage'],
            'matched_terms': result['matched_terms']
        })
    
    if len(search_results) > 5:
        answer_parts.append(f"\n...aur {len(search_results) - 5} documents mein bhi related content mila.")
    
    response['answer'] = '\n'.join(answer_parts)
    response['suggestions'] = _generate_suggestions(query_lower)
    
    return response


def _generate_no_result_answer(query_lower):
    """Generate helpful response when no documents match"""
    suggestions = []
    
    # Check broad topics
    if any(w in query_lower for w in ['solder', 'string']):
        suggestions.append("'Soldering' ya 'Stringer' se related documents upload karein")
    elif any(w in query_lower for w in ['lamin']):
        suggestions.append("'Lamination' process ke documents upload karein")
    elif any(w in query_lower for w in ['pack', 'dispatch']):
        suggestions.append("'Packing & Dispatch' documents upload karein")
    elif any(w in query_lower for w in ['test', 'flash', 'el ']):
        suggestions.append("'Testing & QC' documents upload karein")
    
    answer = "❌ Is query se related koi document content nahi mila.\n\n"
    answer += "**Possible reasons:**\n"
    answer += "• Abhi tak koi document upload nahi hua ya text extract nahi hua\n"
    answer += "• Documents mein ye specific content available nahi hai\n"
    answer += "• Different keywords try karein\n"
    
    if suggestions:
        answer += f"\n💡 **Suggestion:** {suggestions[0]}"
    
    answer += "\n\n🔄 Documents tab se files upload karein, text automatically extract ho jayega."
    
    return answer


def _generate_suggestions(query_lower):
    """Generate follow-up suggestions based on query"""
    suggestions = []
    
    topic_map = {
        'solder': ['ribbon specification', 'flux temperature', 'stringer setting'],
        'lamin': ['lamination temperature profile', 'crosslink test', 'peel test result'],
        'pack': ['packing checklist', 'label placement', 'pallet specification'],
        'test': ['flash test parameters', 'EL inspection criteria', 'hi-pot voltage'],
        'frame': ['frame torque value', 'sealant application', 'corner key specification'],
        'junction': ['junction box specification', 'potting procedure', 'diode test'],
        'cell': ['cell specification', 'cell grading criteria', 'incoming inspection'],
        'calibr': ['calibration schedule', 'calibration certificate', 'measurement uncertainty'],
        'audit': ['audit checklist', 'audit finding', 'corrective action'],
        'train': ['training matrix', 'competency record', 'skill assessment'],
        'iso': ['ISO 9001 clause', 'management review', 'quality objective'],
        'ncr': ['NCR status', 'root cause analysis', 'corrective action'],
        'spec': ['material specification', 'product specification', 'process parameter'],
    }
    
    for key, suggs in topic_map.items():
        if key in query_lower:
            suggestions.extend(suggs[:3])
            break
    
    if not suggestions:
        suggestions = ['inspection criteria', 'process parameter', 'specification']
    
    return suggestions[:4]
